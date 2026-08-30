import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "PathCompanion AI — An Agentic Career and Talent Companion System"
subtitle.text = ("Development Project-CIS6035 /Dissertation Project CIS 6000\n"
                 "BSc. Software Engineering\n"
                 "Student: Sainthavy Yogaraththinam\n"
                 "Supervisor: [Supervisor Name]\n"
                 "ICBT Campus, Jaffna\n"
                 "Date: ")

# Slide 2: Table of Contents
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Table of Contents"
content = slide.placeholders[1]
content.text = ("Introduction\n"
                "Problem Statement\n"
                "Objectives\n"
                "Literature Review\n"
                "Proposed Solution\n"
                "Functionalities & Non-functionalities\n"
                "Methodology\n"
                "Technology Stack\n"
                "Limitations\n"
                "Lesson Learned\n"
                "Future Recommendation\n"
                "References")

# Slide 3: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Introduction"
content = slide.placeholders[1]
content.text = ("Career development is fragmented across disconnected platforms.\n"
                "PathCompanion AI provides a single full-stack web application.\n"
                "Integrates six cooperating AI agents through a shared memory and event bus.\n"
                "Hybrid intelligence design using two custom-trained machine-learning models.")

# Slide 4: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Problem Statement"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Career development today is fragmented across disconnected tools."
p = tf.add_paragraph()
p.text = "Job discovery is inefficient (semantic mismatch)."
p.level = 1
p = tf.add_paragraph()
p.text = "Skill gaps are invisible without structured feedback."
p.level = 1
p = tf.add_paragraph()
p.text = "Learning paths are unguided."
p.level = 1
p = tf.add_paragraph()
p.text = "Resumes are generic and rarely ATS-optimised per job."
p.level = 1
p = tf.add_paragraph()
p.text = "Interview anxiety is unaddressed without realistic rehearsal."
p.level = 1

# Slide 5: Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Objectives"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Main Objective: Design, implement, and evaluate a multi-agent AI web application for career support."
p = tf.add_paragraph()
p.text = "Train and evaluate a resume-category classification model."
p.level = 1
p = tf.add_paragraph()
p.text = "Train and evaluate a multi-label skill-extraction model."
p.level = 1
p = tf.add_paragraph()
p.text = "Implement automated job aggregation, gap analysis, learning plans, and ATS resume generation."
p.level = 1
p = tf.add_paragraph()
p.text = "Conduct realistic voice-based mock interviews."
p.level = 1
p = tf.add_paragraph()
p.text = "Operate within free-tier limits."
p.level = 1

# Slide 6: Literature Review
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Literature Review"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Key Areas Researched:"
p = tf.add_paragraph()
p.text = "Fragmentation of Career Tooling: Gap in integration among platforms."
p.level = 1
p = tf.add_paragraph()
p.text = "Agentic AI Systems: Using event-driven multi-agent architectures."
p.level = 1
p = tf.add_paragraph()
p.text = "Text Classification: Transfer learning with DistilBERT over TF-IDF baselines."
p.level = 1
p = tf.add_paragraph()
p.text = "Skill Extraction: LLM-assisted weak supervision for labelling."
p.level = 1
p = tf.add_paragraph()
p.text = "Semantic Retrieval: Using pgvector for dense embedding search."
p.level = 1

# Slide 7: Proposed Solution
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Proposed Solution"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Web application with React (Vite), FastAPI, and Supabase."
p = tf.add_paragraph()
p.text = "Six Coordinated Agents:"
p.level = 1
p = tf.add_paragraph()
p.text = "Job Scout, Skill Gap Analyzer, Learning Resource Finder"
p.level = 2
p = tf.add_paragraph()
p.text = "Resume Architect, Mock Interview Agent, Career Path Planner"
p.level = 2
p = tf.add_paragraph()
p.text = "Interconnected Cascade: Actions in one agent (e.g. learning a skill) automatically update states in others (e.g. gap analysis, resume drafts)."
p.level = 1

# Slide 8: Functionalities & Non-functionalities
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Functionalities & Non-functionalities"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Functional Requirements:"
p = tf.add_paragraph()
p.text = "Job aggregation & semantic matching."
p.level = 1
p = tf.add_paragraph()
p.text = "Skill extraction & gap analysis."
p.level = 1
p = tf.add_paragraph()
p.text = "Tailored ATS resume generation."
p.level = 1
p = tf.add_paragraph()
p.text = "Voice mock interviews with panic control."
p.level = 1
p = tf.add_paragraph()
p.text = "Non-Functional Requirements:"
p.level = 0
p = tf.add_paragraph()
p.text = "Security (Row-Level Security in Supabase)."
p.level = 1
p = tf.add_paragraph()
p.text = "Performance (sub-2s latency for interviews)."
p.level = 1
p = tf.add_paragraph()
p.text = "Zero-cost operation (Free tiers)."
p.level = 1

# Slide 9: Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Methodology"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Agile Life Cycle:"
p = tf.add_paragraph()
p.text = "16 weeks divided into 8 two-week sprints."
p.level = 1
p = tf.add_paragraph()
p.text = "Iterative development to handle ML experiments & API integration."
p.level = 1
p = tf.add_paragraph()
p.text = "Requirement Gathering:"
p.level = 0
p = tf.add_paragraph()
p.text = "Questionnaires (quantitative prevalence of problems)."
p.level = 1
p = tf.add_paragraph()
p.text = "Semi-structured interviews (qualitative pain points)."
p.level = 1
p = tf.add_paragraph()
p.text = "Platform analysis of existing tools."
p.level = 1

# Slide 10: Technology Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Technology Stack"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Frontend: React 18, Vite, Vanilla CSS."
p = tf.add_paragraph()
p.text = "Backend: FastAPI (Python 3.11, async), Uvicorn."
p.level = 1
p = tf.add_paragraph()
p.text = "Database: Supabase (PostgreSQL, pgvector, HNSW indexes)."
p.level = 1
p = tf.add_paragraph()
p.text = "Machine Learning: scikit-learn, Hugging Face Transformers, PyTorch."
p.level = 1
p = tf.add_paragraph()
p.text = "AI Services:"
p.level = 1
p = tf.add_paragraph()
p.text = "Google Gemini Flash (generation/embeddings)"
p.level = 2
p = tf.add_paragraph()
p.text = "Groq Llama 3.3 70B & Whisper (chat/STT)"
p.level = 2

# Slide 11: Limitations
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Limitations"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Model Truncation: Macro-F1 below target for Model 1 due to 256-token truncation of long resumes."
p = tf.add_paragraph()
p.text = "Free-Tier Limits: System is dependent on external LLM and API free-tier quotas."
p.level = 0
p = tf.add_paragraph()
p.text = "Resource Constraints: Single developer limited parallel workstreams."
p.level = 0
p = tf.add_paragraph()
p.text = "Language: English-only interface in this release."
p.level = 0

# Slide 12: Lesson Learned
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Lesson Learned"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Hybrid AI approach effectively reduces external dependencies and costs."
p = tf.add_paragraph()
p.text = "Managing free-tier rate limits requires careful caching and background processing."
p.level = 0
p = tf.add_paragraph()
p.text = "Continuous user feedback is vital for relevance."
p.level = 0
p = tf.add_paragraph()
p.text = "Event-driven architecture ensures agents are loosely coupled."
p.level = 0

# Slide 13: Future Recommendation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Future Recommendation"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Real-world Release: Deploy system to students and early-career professionals."
p = tf.add_paragraph()
p.text = "Recruiter Portal: Add features for employers."
p.level = 0
p = tf.add_paragraph()
p.text = "Multilingual Support: Expand to Tamil and Sinhala."
p.level = 0
p = tf.add_paragraph()
p.text = "Mobile Application: Develop a mobile client on the same backend."
p.level = 0

# Slide 14: References
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "References"
content = slide.placeholders[1]
tf = content.text_frame
tf.text = "Sanh et al. (2019). DistilBERT, a distilled version of BERT: smaller, faster, cheaper and lighter."
p = tf.add_paragraph()
p.text = "Pedregosa et al. (2011). Scikit-learn: Machine Learning in Python."
p.level = 0
p = tf.add_paragraph()
p.text = "Zhang, Jensen and Plank (2022). Weak Supervision for text classification."
p.level = 0
p = tf.add_paragraph()
p.text = "Supabase (2026). pgvector documentation."
p.level = 0

prs.save("PathCompanionAI_Final_Presentation.pptx")
print("Presentation saved successfully.")
