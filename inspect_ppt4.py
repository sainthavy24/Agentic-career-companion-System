from pptx import Presentation

prs = Presentation(r"C:\Users\asus\Documents\GitHub\Agentic-career-companion-System\docs\PathCompanionAI_PRES1_Presentation_v2.pptx")
shape = prs.slides[0].shapes[0]
if shape.has_text_frame:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            print(f"Font name: {run.font.name}, size: {run.font.size}, color: {run.font.color.rgb if run.font.color and hasattr(run.font.color, 'rgb') else 'Not set'}")
