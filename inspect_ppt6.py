from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

prs = Presentation(r"C:\Users\asus\Documents\GitHub\Agentic-career-companion-System\docs\PathCompanionAI_PRES1_Presentation_v2.pptx")
slide = prs.slides[0]
for j, shape in enumerate(slide.shapes):
    print(f"Shape {j} type: {shape.shape_type}")
