from pptx import Presentation

prs = Presentation(r"C:\Users\asus\Documents\GitHub\Agentic-career-companion-System\docs\PathCompanionAI_PRES1_Presentation_v2.pptx")
slide = prs.slides[0]
for j, shape in enumerate(slide.shapes):
    print(f"Shape {j} type: {shape.shape_type}")
    if hasattr(shape, "width"):
        print(f"  Width: {shape.width}, Height: {shape.height}")
    if hasattr(shape, "text_frame"):
        print(f"  Text: {shape.text[:30]}")
