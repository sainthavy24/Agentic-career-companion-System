import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

slides = []
for _ in range(14):
    slide = prs.slides.add_slide(blank_layout)
    # Professional dark blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 32, 67)
    bg.line.fill.background()
    slides.append(slide)

def add_title(slide, text):
    left = Inches(0.8)
    top = Inches(0.5)
    width = Inches(11.7)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.name = 'Segoe UI'
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Add an accent line under the title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 162, 255)
    line.line.fill.background()
    
def add_body(slide, points, top_in=1.9):
    left = Inches(0.8)
    top = Inches(top_in)
    width = Inches(11.7)
    height = Inches(5.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, (text, level) in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = text
        p.level = level
        p.font.size = Pt(28 - (level * 4)) # 28 for level 0, 24 for level 1
        p.font.name = 'Segoe UI'
        p.font.color.rgb = RGBColor(235, 240, 245)
        p.space_after = Pt(14)
        if level == 0:
            p.space_before = Pt(10)

# Slide 1: Title
slide = slides[0]
left = Inches(1)
top = Inches(2.2)
width = Inches(11.3)
height = Inches(3)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PathCompanion AI"
p.font.bold = True
p.font.size = Pt(64)
p.font.name = 'Segoe UI'
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "An Agentic Career and Talent Companion System"
p2.font.size = Pt(36)
p2.font.name = 'Segoe UI'
p2.font.color.rgb = RGBColor(0, 162, 255)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(15)

p3 = tf.add_paragraph()
p3.text = "\nBSc. Software Engineering\nStudent: Sainthavy Yogaraththinam\nSupervisor: [Supervisor Name]\nICBT Campus, Jaffna"
p3.font.size = Pt(22)
p3.font.name = 'Segoe UI'
p3.font.color.rgb = RGBColor(200, 200, 200)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(20)

# Slide 2: Table of Contents
add_title(slides[1], "Table of Contents")
add_body(slides[1], [
    ("Introduction", 0),
    ("Problem Statement", 0),
    ("Objectives", 0),
    ("Literature Review", 0),
    ("Proposed Solution", 0),
    ("Functionalities & Non-functionalities", 0),
    ("Methodology & Technology Stack", 0),
    ("Limitations & Lessons Learned", 0),
    ("Future Recommendation & References", 0)
])

# Slide 3: Introduction
add_title(slides[2], "Introduction")
add_body(slides[2], [
    ("Career development today is fragmented across disconnected platforms.", 0),
    ("PathCompanion AI provides a single full-stack web application.", 0),
    ("Integrates six cooperating AI agents through a shared memory and event bus:", 0),
    ("Job Scout, Skill Gap Analyzer, Learning Resource Finder", 1),
    ("Resume Architect, Mock Interview Agent, Career Path Planner", 1),
    ("Hybrid intelligence design using two custom-trained machine-learning models.", 0)
])

# Slide 4: Problem Statement
add_title(slides[3], "Problem Statement")
add_body(slides[3], [
    ("Career development today is fragmented across disconnected tools.", 0),
    ("Job discovery is inefficient due to semantic mismatch.", 1),
    ("Skill gaps are invisible without structured feedback.", 1),
    ("Learning paths are unguided and disconnected from the gaps.", 1),
    ("Resumes are generic and rarely ATS-optimised per job.", 1),
    ("Interview anxiety is unaddressed without realistic rehearsal.", 1),
    ("Career progression timelines and next steps are opaque.", 1)
])

# Slide 5: Objectives
add_title(slides[4], "Objectives")
add_body(slides[4], [
    ("Main Objective", 0),
    ("Design, implement, and evaluate a multi-agent AI web application for career support.", 1),
    ("Sub-Objectives", 0),
    ("Train and evaluate a resume-category classification model.", 1),
    ("Train and evaluate a multi-label skill-extraction model.", 1),
    ("Implement automated job aggregation, gap analysis, and ATS resume generation.", 1),
    ("Conduct realistic voice-based mock interviews.", 1),
    ("Operate entirely within free-tier limits.", 1)
])

# Slide 6: Literature Review
add_title(slides[5], "Literature Review")
add_body(slides[5], [
    ("Fragmentation of Career Tooling", 0),
    ("Major gap in integration among existing platforms.", 1),
    ("Agentic AI Systems", 0),
    ("Effectiveness of event-driven multi-agent architectures.", 1),
    ("Text Classification", 0),
    ("Transfer learning with DistilBERT over TF-IDF baselines.", 1),
    ("Skill Extraction", 0),
    ("LLM-assisted weak supervision for labelling job descriptions.", 1),
    ("Semantic Retrieval", 0),
    ("Using pgvector for dense embedding search in PostgreSQL.", 1)
])

# Slide 7: Proposed Solution
add_title(slides[6], "Proposed Solution")
add_body(slides[6], [
    ("Web application (React, FastAPI, Supabase)", 0),
    ("Six Coordinated Agents:", 0),
    ("Job Scout, Skill Gap Analyzer, Learning Resource Finder", 1),
    ("Resume Architect, Mock Interview Agent, Career Path Planner", 1),
    ("Interconnected Cascade System:", 0),
    ("Actions in one agent automatically update states in others.", 1),
    ("Example: Learning a skill automatically refreshes gap analysis and resume drafts.", 1)
])

# Slide 8: Functionalities
add_title(slides[7], "Functionalities & Non-functionalities")
add_body(slides[7], [
    ("Functional Requirements:", 0),
    ("Job aggregation & semantic matching.", 1),
    ("Skill extraction & automated gap analysis.", 1),
    ("Tailored ATS resume generation.", 1),
    ("Voice mock interviews with panic control.", 1),
    ("Non-Functional Requirements:", 0),
    ("Security (Row-Level Security in Supabase).", 1),
    ("Performance (sub-2s latency for interviews).", 1),
    ("Zero-cost operation (Free tiers only).", 1)
])

# Slide 9: Methodology
add_title(slides[8], "Methodology")
add_body(slides[8], [
    ("Agile Life Cycle:", 0),
    ("16 weeks divided into 8 two-week sprints.", 1),
    ("Iterative development to handle ML experiments & API integration.", 1),
    ("Requirement Gathering:", 0),
    ("Questionnaires to evaluate quantitative prevalence of problems.", 1),
    ("Semi-structured interviews to understand qualitative pain points.", 1),
    ("Platform analysis of existing disconnected tools.", 1)
])

# Slide 10: Technology Stack
add_title(slides[9], "Technology Stack")
add_body(slides[9], [
    ("Frontend: React 18, Vite, Vanilla CSS.", 0),
    ("Backend: FastAPI (Python 3.11, async), Uvicorn.", 0),
    ("Database: Supabase (PostgreSQL, pgvector, HNSW indexes).", 0),
    ("Machine Learning: scikit-learn, Hugging Face Transformers, PyTorch.", 0),
    ("AI Services:", 0),
    ("Google Gemini Flash (generation/embeddings).", 1),
    ("Groq Llama 3.3 70B & Whisper (chat/STT).", 1)
])

# Slide 11: Limitations
add_title(slides[10], "Limitations")
add_body(slides[10], [
    ("Model Truncation:", 0),
    ("Macro-F1 below target for Model 1 due to 256-token truncation of long resumes.", 1),
    ("Free-Tier Limits:", 0),
    ("System is entirely dependent on external LLM and API free-tier quotas.", 1),
    ("Resource Constraints:", 0),
    ("Single developer limited parallel workstreams.", 1),
    ("Language:", 0),
    ("English-only interface in this release.", 1)
])

# Slide 12: Lesson Learned
add_title(slides[11], "Lesson Learned")
add_body(slides[11], [
    ("Hybrid AI Architecture:", 0),
    ("Effectively reduces external dependencies and operational costs.", 1),
    ("Rate Limit Management:", 0),
    ("Managing free-tier rate limits requires careful caching and background processing.", 1),
    ("User Centric Iteration:", 0),
    ("Continuous user feedback (e.g., adding mock-exams) is vital for relevance.", 1),
    ("System Design:", 0),
    ("Event-driven architecture ensures agents are loosely coupled and robust.", 1)
])

# Slide 13: Future Recommendation
add_title(slides[12], "Future Recommendation")
add_body(slides[12], [
    ("Real-world Release:", 0),
    ("Deploy system to students and early-career professionals.", 1),
    ("Recruiter Portal:", 0),
    ("Add features allowing employers to find candidates.", 1),
    ("Multilingual Support:", 0),
    ("Expand to Tamil and Sinhala for broader accessibility.", 1),
    ("Mobile Application:", 0),
    ("Develop a mobile client using the same unified backend.", 1)
])

# Slide 14: References
add_title(slides[13], "References")
add_body(slides[13], [
    ("Sanh et al. (2019). DistilBERT, a distilled version of BERT.", 0),
    ("Pedregosa et al. (2011). Scikit-learn: Machine Learning in Python.", 0),
    ("Zhang, Jensen and Plank (2022). Weak Supervision for text classification.", 0),
    ("Supabase (2026). pgvector documentation.", 0)
])

prs.save("PathCompanionAI_Final_Presentation_Perfect.pptx")
print("Perfect Presentation saved successfully.")
