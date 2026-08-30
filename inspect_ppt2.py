from pptx import Presentation

prs = Presentation(r"C:\Users\asus\Documents\GitHub\Agentic-career-companion-System\docs\PathCompanionAI_PRES1_Presentation_v2.pptx")
for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        if hasattr(shape, "text"):
            text = shape.text[:50].replace('\n', ' ')
            print(f"Shape {j}: {text}...")
